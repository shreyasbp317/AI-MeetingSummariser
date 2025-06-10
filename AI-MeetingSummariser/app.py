from flask import Flask, render_template, request, jsonify
from summarizer import generate_summary
from pptx import Presentation
import os
import logging

app = Flask(__name__)

# Configure logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

# Upload configuration
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB limit


def extract_ppt_text(filepath):
    """Enhanced PPT text extraction"""
    try:
        prs = Presentation(filepath)
        full_text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        full_text.append(paragraph.text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            full_text.append(cell.text)
        return '\n'.join([t for t in full_text if t.strip()])
    except Exception as e:
        logger.error(f"PPT Extraction Error: {str(e)}")
        raise


@app.route('/')
def home():
    return render_template('index.html')


@app.route('/summarize-ppt', methods=['POST'])
def summarize_ppt():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400

        file = request.files['file']
        if not file or file.filename == '':
            return jsonify({'error': 'No selected file'}), 400

        if not file.filename.lower().endswith('.pptx'):
            return jsonify({'error': 'Only .pptx files are supported'}), 400

        # Save file
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(filepath)
        logger.info(f"File saved to {filepath}")

        # Extract and summarize
        ppt_text = extract_ppt_text(filepath)
        logger.info(f"Extracted text length: {len(ppt_text)} chars")

        if not ppt_text:
            return jsonify({'error': 'No readable text found in PPT'}), 400

        summary = generate_summary(ppt_text)
        os.remove(filepath)  # Clean up

        return jsonify({'summary': summary})

    except Exception as e:
        logger.error(f"Summarization error: {str(e)}")
        return jsonify({'error': f'Processing failed: {str(e)}'}), 500


if __name__ == '__main__':
    app.run(debug=True, port=5000)